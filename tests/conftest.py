from __future__ import annotations

from pathlib import Path

import pytest
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas

_CJK_FONT = "HeiseiMin-W3"
pdfmetrics.registerFont(UnicodeCIDFont(_CJK_FONT))


def _write_pdf(path: Path, pages: list[list[str]]) -> None:
    c = canvas.Canvas(str(path), pagesize=A4)
    for page_index, lines in enumerate(pages):
        c.setFont(_CJK_FONT, 12)
        y = 800
        for line in lines:
            c.drawString(72, y, line)
            y -= 14
        if page_index < len(pages) - 1:
            c.showPage()
    c.save()


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    """A 2-page text PDF with distinct, searchable content per page."""
    pdf_path = tmp_path / "sample.pdf"
    _write_pdf(pdf_path, [
        ["東京の天気予報について", "本日は晴れ、最高気温は25度です。", "!!!注意!!!"],
        ["大阪の天気予報について", "本日は雨、最高気温は18度です。"],
    ])
    return pdf_path


@pytest.fixture
def blank_pdf(tmp_path: Path) -> Path:
    """A PDF with a page that has no text (simulates a scanned page)."""
    pdf_path = tmp_path / "blank.pdf"
    _write_pdf(pdf_path, [[]])
    return pdf_path


def _write_pptx(path: Path, slides: list[list[str]]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for lines in slides:
        slide = prs.slides.add_slide(blank_layout)
        if lines:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            tf = txBox.text_frame
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
    prs.save(str(path))


def _write_xlsx(path: Path, sheets: dict[str, list[list[str]]]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    for i, (name, rows) in enumerate(sheets.items()):
        ws = wb.active if i == 0 else wb.create_sheet()
        ws.title = name
        for row in rows:
            ws.append(row)
    wb.save(str(path))


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """A 2-slide pptx with distinct, searchable content per slide."""
    path = tmp_path / "sample.pptx"
    _write_pptx(path, [
        ["東京の天気予報について", "本日は晴れ、最高気温は25度です。"],
        ["大阪の天気予報について", "本日は雨、最高気温は18度です。"],
    ])
    return path


@pytest.fixture
def blank_pptx(tmp_path: Path) -> Path:
    """A pptx with a single blank slide (no text shapes)."""
    path = tmp_path / "blank.pptx"
    _write_pptx(path, [[]])
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    """A 2-sheet xlsx with distinct, searchable content per sheet."""
    path = tmp_path / "sample.xlsx"
    _write_xlsx(path, {
        "東京": [
            ["東京の天気予報について"],
            ["本日は晴れ、最高気温は25度です。"],
        ],
        "大阪": [
            ["大阪の天気予報について"],
            ["本日は雨、最高気温は18度です。"],
        ],
    })
    return path


@pytest.fixture
def blank_xlsx(tmp_path: Path) -> Path:
    """An xlsx with a single empty sheet."""
    path = tmp_path / "blank.xlsx"
    _write_xlsx(path, {"Sheet1": []})
    return path
