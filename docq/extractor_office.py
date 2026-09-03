"""Office document (pptx / xlsx) -> per-page plain text extraction.

Each slide (pptx) or sheet (xlsx) maps to one "page" of text, following the
same interface as extractor_pdf.
"""
from __future__ import annotations

from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation


def extract_pages(path: Path) -> list[str]:
    """Return a list of page texts from an Office document.

    For .pptx: one entry per slide.
    For .xlsx: one entry per sheet.
    A slide/sheet with no text yields "".
    """
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    raise ValueError(f"Unsupported format: {suffix}")


def _extract_pptx(path: Path) -> list[str]:
    prs = Presentation(str(path))
    pages: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        pages.append("\n".join(texts))
    return pages


def _extract_xlsx(path: Path) -> list[str]:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[str] = []
    try:
        for ws in wb.worksheets:
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                line = " ".join(cells)
                if line:
                    rows.append(line)
            pages.append("\n".join(rows))
    finally:
        wb.close()
    return pages
